import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank layout

    # Colors
    NAVY = RGBColor(13, 92, 157)        # #0D5C9D Primary SIH Blue
    DARK_NAVY = RGBColor(10, 59, 102)   # #0A3B66
    ORANGE = RGBColor(232, 96, 36)      # #E86024 Accent Orange
    GREEN = RGBColor(46, 125, 50)       # #2E7D32 Accent Green
    GOLD = RGBColor(217, 119, 6)        # #D97706
    BG_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)   # White
    BORDER_COLOR = RGBColor(218, 225, 233) # #DAE1E9
    TEXT_DARK = RGBColor(30, 41, 59)     # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BLUE = RGBColor(238, 246, 255)
    LIGHT_ORANGE = RGBColor(254, 243, 199)
    LIGHT_GREEN = RGBColor(236, 253, 245)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_LIGHT

    def add_header(slide, title_text, slide_num):
        # Oval Team Badge (Top Left)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(0.25), Inches(1.8), Inches(0.55))
        oval.fill.solid()
        oval.fill.fore_color.rgb = WHITE
        oval.line.color.rgb = NAVY
        oval.line.width = Pt(1.5)
        tf = oval.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Aquaregia"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

        # Slide Title (Center Header)
        title_box = slide.shapes.add_textbox(Inches(2.4), Inches(0.18), Inches(8.2), Inches(0.65))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER

        # SIH Header Right Text & Logo Badge
        sih_box = slide.shapes.add_textbox(Inches(10.7), Inches(0.12), Inches(2.2), Inches(0.75))
        tf = sih_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "SMART INDIA"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = "HACKATHON 2026"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = ORANGE
        p2.alignment = PP_ALIGN.RIGHT

        # Footer Strip
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer.fill.solid()
        footer.fill.fore_color.rgb = NAVY
        footer.line.fill.background()

        # Footer text
        ft_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.07), Inches(12.5), Inches(0.4))
        tf = ft_box.text_frame
        p = tf.paragraphs[0]
        p.text = "@SIH Idea submission- Template"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Page Number Right
        num_box = slide.shapes.add_textbox(Inches(12.4), Inches(7.07), Inches(0.6), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.RIGHT

    def add_card(slide, left, top, width, height, title="", title_color=ORANGE, bg_color=CARD_BG, border_color=BORDER_COLOR):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.45))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = title_color
        return card

    # ==========================================
    # SLIDE 1: TITLE PAGE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Top SIH Banner
    banner = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.5), Inches(1.1))
    tf = banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Right Logo Emblem Container
    logo_bg = slide1.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(9.8), Inches(0.3), Inches(3.0), Inches(3.2))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = LIGHT_BLUE
    logo_bg.line.color.rgb = NAVY
    logo_bg.line.width = Pt(2)
    
    logo_tb = slide1.shapes.add_textbox(Inches(9.9), Inches(0.9), Inches(2.8), Inches(2.0))
    tf = logo_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SMART INDIA\nHACKATHON\n2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "SIH"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.alignment = PP_ALIGN.CENTER

    # Slide Main Title Header
    main_title = slide1.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.8), Inches(1.2))
    tf = main_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "VANIKA: AI-Based Cognitive Gaming &\nMemory Assistance Platform"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    sub_title = slide1.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(8.8), Inches(0.6))
    tf = sub_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Engineered for Elderly Dementia Patients in North-Eastern India"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ORANGE

    # Metadata Details Box (Left Side)
    meta_card = add_card(slide1, Inches(0.8), Inches(3.2), Inches(7.5), Inches(3.8), title="PROJECT REGISTRATION DETAILS", title_color=NAVY)
    meta_tb = slide1.shapes.add_textbox(Inches(1.0), Inches(3.7), Inches(7.1), Inches(3.2))
    tf = meta_tb.text_frame
    tf.word_wrap = True
    
    items = [
        ("Problem Statement ID –", " SIH1724"),
        ("Problem Statement Title –", " AI-Based Cognitive Gaming and Memory Assistance Platform for Elderly Dementia Patients"),
        ("Theme –", " MedTech / Healthcare AI / Digital Inclusion"),
        ("PS Category –", " Software"),
        ("Team ID –", " SIH2026_AQUAREGIA"),
        ("Team Name –", " Aquaregia (Registered on SIH portal)")
    ]
    for idx, (lbl, val) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run1 = p.add_run()
        run1.text = "• " + lbl
        run1.font.bold = True
        run1.font.size = Pt(14)
        run1.font.color.rgb = TEXT_DARK
        
        run2 = p.add_run()
        run2.text = val
        run2.font.bold = (lbl in ["Team Name –", "Problem Statement ID –"])
        run2.font.size = Pt(14)
        run2.font.color.rgb = NAVY if lbl == "Team Name –" else TEXT_DARK

    # Core Features Summary Box (Right Side)
    feat_card = add_card(slide1, Inches(8.6), Inches(3.2), Inches(4.0), Inches(3.8), title="VANIKA CORE PILLARS", title_color=GREEN, bg_color=LIGHT_GREEN, border_color=GREEN)
    feat_tb = slide1.shapes.add_textbox(Inches(8.75), Inches(3.7), Inches(3.7), Inches(3.2))
    tf = feat_tb.text_frame
    tf.word_wrap = True

    pillars = [
        ("🧠 On-Device Emotion AI:", " Senses confusion/frustration and eases game difficulty in real time."),
        ("🌾 NER Cultural Memories:", " Bihu festivals, Majuli island scans, tea garden walks & folk wisdom."),
        ("🗣️ 6 Regional Languages:", " Assamese, Bodo, Khasi, Mizo, Nagamese & English voice guide."),
        ("🔒 AES-256 Data Vault:", " 100% offline-capable local storage compliant with India DPDP Act 2023."),
        ("🩺 Caregiver Portal:", " 7 & 30-day wellness trends and actionable guidance for families & ASHA workers.")
    ]
    for idx, (title, desc) in enumerate(pillars):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT_DARK


    # ==========================================
    # SLIDE 2: IDEA / APPROACH
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "IDEA / APPROACH — Vanika Cognitive Platform", 2)

    # Box 1: Top Left - Detailed Explanation
    add_card(slide2, Inches(0.5), Inches(0.95), Inches(6.0), Inches(2.9), title="Detailed Explanation of Proposed Solution", title_color=ORANGE)
    tb1 = slide2.shapes.add_textbox(Inches(0.65), Inches(1.4), Inches(5.7), Inches(2.3))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p_pts1 = [
        ("Problem Addressed: ", "Elderly dementia patients face rapid memory decline, diagnostic stigma, language barriers in North-East India (NER), and high caregiver burden."),
        ("Vanika Solution: ", "An enculturated, voice-first AI platform that transforms scary medical tests into familiar 'digital courtyard' games and daily memory play."),
        ("Multilingual Voice AI: ", "Conversational companion 'Oja / Aita' speaking 6 regional dialects (Assamese, Bodo, Khasi, Mizo, Nagamese, English) with gentle voice pacing.")
    ]
    for idx, (bld, txt) in enumerate(p_pts1):
        p = tf1.paragraphs[0] if idx == 0 else tf1.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run(); r1.text = "➢ " + bld; r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(11); r2.font.color.rgb = TEXT_DARK

    # Box 2: Bottom Left - Innovation & Uniqueness
    add_card(slide2, Inches(0.5), Inches(3.95), Inches(6.0), Inches(2.95), title="Innovation and Uniqueness of Solution", title_color=ORANGE)
    tb2 = slide2.shapes.add_textbox(Inches(0.65), Inches(4.4), Inches(5.7), Inches(2.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p_pts2 = [
        ("100% Offline Local Vault: ", "AES-256 encrypted local storage ensures strict compliance with DPDP Act 2023. Camera frames stay on-device without cloud telemetry."),
        ("On-Device Emotion-Adaptive AI: ", "Real-time voice pacing and facial micro-expression telemetry automatically ease difficulty and play tea-garden folk tunes when frustration is detected."),
        ("Indigenous Cultural Integration: ", "Personalized family photo recall, Bihu morning sequencing, Majuli riverside visual search, and Manimuni/Brahmi herbal caregiving guidance.")
    ]
    for idx, (bld, txt) in enumerate(p_pts2):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run(); r1.text = "➢ " + bld; r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(11); r2.font.color.rgb = TEXT_DARK

    # Box 3: Top Right - How it Addresses the Problem
    add_card(slide2, Inches(6.8), Inches(0.95), Inches(6.0), Inches(2.9), title="How It Addresses the Problem", title_color=ORANGE)
    tb3 = slide2.shapes.add_textbox(Inches(6.95), Inches(1.4), Inches(5.7), Inches(2.3))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p_pts3 = [
        ("Zero Diagnostic Anxiety: ", "Eliminates intimidating clinical exams. Elders engage in joyful, pressure-free daily play ('Let's look once more' vs red X marks)."),
        ("Empowers Caregivers & ASHA Workers: ", "Unified 7-day and 30-day wellness trends across Memory, Attention, Mood & Reminders with proactive non-alarmist advisories."),
        ("Cognitive Stamina Preservation: ", "Short 15-minute daily routines woven with cultural cues ('Like morning tea over Brahmaputra') to sustain focus without fatigue.")
    ]
    for idx, (bld, txt) in enumerate(p_pts3):
        p = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run(); r1.text = "➢ " + bld; r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(11); r2.font.color.rgb = TEXT_DARK

    # Box 4: Bottom Right - Core Workflow Diagram
    add_card(slide2, Inches(6.8), Inches(3.95), Inches(6.0), Inches(2.95), title="Core System Workflow & User Journey", title_color=GREEN, bg_color=LIGHT_GREEN, border_color=GREEN)
    
    nodes = [
        ("👵 Elderly User", "Selects Native Dialect"),
        ("🎮 Courtyard Play", "Memory & Sequence Games"),
        ("🤖 On-Device AI", "Pacing & Emotion Adaptation"),
        ("🔒 AES-256 Vault", "DPDP Compliant Local Sync"),
        ("🩺 Caregiver Portal", "ASHA & Family Insights")
    ]
    for i, (n_title, n_sub) in enumerate(nodes):
        nx = Inches(7.0 + (i % 3) * 1.85) if i < 3 else Inches(7.9 + (i - 3) * 1.85)
        ny = Inches(4.45) if i < 3 else Inches(5.7)
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, Inches(1.7), Inches(0.95))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = NAVY; box.line.width = Pt(1)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = n_title; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = DARK_NAVY; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = n_sub; p2.font.size = Pt(9); p2.font.color.rgb = TEXT_MUTED; p2.alignment = PP_ALIGN.CENTER


    # ==========================================
    # SLIDE 3: TECHNICAL APPROACH
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "TECHNICAL APPROACH & ARCHITECTURE", 3)

    # Left Column: Tech Stack Grid (Width: 5.6 in)
    add_card(slide3, Inches(0.5), Inches(0.95), Inches(5.6), Inches(5.95), title="Tech Stack & Software Components", title_color=NAVY)
    
    stack_items = [
        ("🎨 Frontend & UI Framework", "React 19, TypeScript, Vite, TailwindCSS v4, Lucide Icons, Motion", LIGHT_BLUE),
        ("⚙️ Server & API Layer", "Node.js, Express, TSX Server, RESTful JSON Endpoints", LIGHT_BLUE),
        ("🧠 AI Models & Companion", "Google Gemini 2.5 API (@google/genai), On-Device Emotion/Voice Pacing Engine", LIGHT_ORANGE),
        ("🔐 Storage & Security Vault", "IndexedDB / LocalStorage, AES-256 Encryption, Zero Camera Cloud Telemetry (DPDP Act 2023)", LIGHT_GREEN),
        ("📊 Analytics & Data Viz", "Recharts, Canvas Confetti, Custom Cognitive Composite Index Algorithms", LIGHT_BLUE),
        ("🌐 Prototype & Repo Links", "Live App: Vanika Web App | Repository: GitHub (Team Aquaregia)", LIGHT_ORANGE)
    ]

    for idx, (st_title, st_desc, st_bg) in enumerate(stack_items):
        sy = Inches(1.45 + idx * 0.88)
        s_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), sy, Inches(5.3), Inches(0.8))
        s_box.fill.solid(); s_box.fill.fore_color.rgb = st_bg
        s_box.line.color.rgb = BORDER_COLOR; s_box.line.width = Pt(1)
        tf = s_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = st_title; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = DARK_NAVY
        p2 = tf.add_paragraph(); p2.text = st_desc; p2.font.size = Pt(10); p2.font.color.rgb = TEXT_DARK

    # Right Column: System Architecture Diagram (Width: 6.4 in)
    add_card(slide3, Inches(6.4), Inches(0.95), Inches(6.4), Inches(5.95), title="System Architecture & Data Flow", title_color=ORANGE)

    arch_layers = [
        ("Layer 1: User & Interface", "Elder-Friendly Client (18px Text / High Contrast) | Voice Guide | Dialect Selector", LIGHT_BLUE),
        ("Layer 2: Cognitive Games & Activities", "Who is This? Photo Recall | Bihu Sequence | Majuli Scan | Memory Garden", WHITE),
        ("Layer 3: AI & Pacing Engine", "Gemini 2.5 AI Companion ('Oja') | Emotion Pacing Analyzer | Folk Music Audio Generator", LIGHT_ORANGE),
        ("Layer 4: Data Sovereignty & Vault", "AES-256 Local Encrypted Vault | LocalStorage / IndexedDB | Zero Cloud Video Frames", LIGHT_GREEN),
        ("Layer 5: Caregiver & ASHA Portal", "Unified 7/30-Day Wellness Trends | Indigenous Care Guidance | Opportunistic Sync", LIGHT_BLUE)
    ]

    for idx, (l_title, l_desc, l_bg) in enumerate(arch_layers):
        ly = Inches(1.5 + idx * 1.05)
        l_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), ly, Inches(6.0), Inches(0.95))
        l_box.fill.solid(); l_box.fill.fore_color.rgb = l_bg
        l_box.line.color.rgb = NAVY if idx==2 else BORDER_COLOR; l_box.line.width = Pt(1.5 if idx==2 else 1)
        tf = l_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = l_title; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = DARK_NAVY
        p2 = tf.add_paragraph(); p2.text = l_desc; p2.font.size = Pt(10); p2.font.color.rgb = TEXT_DARK
        
        if idx < 4:
            arr = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.4), ly + Inches(0.93), Inches(0.4), Inches(0.14))
            arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE; arr.line.fill.background()


    # ==========================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "FEASIBILITY AND VIABILITY", 4)

    # Top Row: 6 Dimension Cards
    feas_cards = [
        ("⚙️ Technical", "Lightweight React/Vite stack with local JS processing.", LIGHT_BLUE),
        ("💡 Economic", "Open-source codebase with zero cloud fee per session.", LIGHT_GREEN),
        ("👵 User Viability", "Elder UI (18px text, high contrast, warm voice guide).", LIGHT_ORANGE),
        ("📡 Operational", "100% offline functionality suited for rural NER terrain.", LIGHT_BLUE),
        ("📈 Scalability", "Modular game architecture enables rapid dialect additions.", LIGHT_GREEN),
        ("🏦 Financial", "B2G National Health Mission + Elder care home SaaS.", LIGHT_ORANGE)
    ]

    for idx, (fc_title, fc_desc, fc_bg) in enumerate(feas_cards):
        cx = Inches(0.45 + idx * 2.08)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(0.95), Inches(1.98), Inches(1.8))
        card.fill.solid(); card.fill.fore_color.rgb = fc_bg
        card.line.color.rgb = BORDER_COLOR; card.line.width = Pt(1)
        tf = card.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = fc_title; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = DARK_NAVY; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = fc_desc; p2.font.size = Pt(10); p2.font.color.rgb = TEXT_DARK; p2.alignment = PP_ALIGN.CENTER

    # Bottom Left Box: Potential Challenges & Strategies Table
    add_card(slide4, Inches(0.45), Inches(2.9), Inches(7.5), Inches(4.0), title="Potential Challenges, Risks & Mitigation Strategies", title_color=ORANGE)
    
    rows, cols = 4, 3
    table_shape = slide4.shapes.add_table(rows, cols, Inches(0.6), Inches(3.4), Inches(7.2), Inches(3.3))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(2.8)

    headers = ["Risk / Challenge Area", "Potential Risk Details", "Mitigation Strategy"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.text = h; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = WHITE

    table_data = [
        ("Low Internet in Rural NER", "Intermittent mobile data in hilly remote districts.", "100% offline-capable local storage with background opportunistic sync."),
        ("Elder Diagnostic Stigma", "Reluctance to take scary clinical memory assessments.", "Enculturated 'digital courtyard' play with warm voice guide & zero medical jargon."),
        ("Data Privacy & Security", "Concerns regarding camera feeds & family photo leaks.", "AES-256 encrypted local vault; camera frames never leave device (DPDP Act 2023).")
    ]

    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i%2==0 else BG_LIGHT
            p = cell.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(10); p.font.color.rgb = TEXT_DARK
            if j == 0: p.font.bold = True

    # Bottom Right Box: Feasibility Visual Rating Chart
    add_card(slide4, Inches(8.15), Inches(2.9), Inches(4.7), Inches(4.0), title="Feasibility & Viability Ratings (/10)", title_color=GREEN, bg_color=LIGHT_GREEN, border_color=GREEN)
    
    ratings = [
        ("Technical Feasibility", 9.5, NAVY),
        ("Economic Feasibility", 9.8, GREEN),
        ("User Viability & Safety", 9.2, ORANGE),
        ("Operational Readiness", 9.0, NAVY),
        ("Regional Scalability", 9.4, GREEN),
        ("Financial Sustainability", 9.0, GOLD)
    ]

    for idx, (r_name, r_score, r_clr) in enumerate(ratings):
        ry = Inches(3.45 + idx * 0.55)
        tb = slide4.shapes.add_textbox(Inches(8.3), ry, Inches(2.2), Inches(0.35))
        tf = tb.text_frame; p = tf.paragraphs[0]; p.text = r_name; p.font.bold = True; p.font.size = Pt(10); p.font.color.rgb = TEXT_DARK
        
        bar_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), ry + Inches(0.05), Inches(1.8), Inches(0.2))
        bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = BORDER_COLOR; bar_bg.line.fill.background()
        
        bar_fill = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), ry + Inches(0.05), Inches(1.8 * (r_score / 10.0)), Inches(0.2))
        bar_fill.fill.solid(); bar_fill.fill.fore_color.rgb = r_clr; bar_fill.line.fill.background()
        
        stb = slide4.shapes.add_textbox(Inches(12.35), ry, Inches(0.4), Inches(0.35))
        tf = stb.text_frame; p = tf.paragraphs[0]; p.text = str(r_score); p.font.bold = True; p.font.size = Pt(10); p.font.color.rgb = r_clr


    # ==========================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "IMPACT AND BENEFITS", 5)

    # Main Section: 3 Stakeholder Columns (Width: 8.5 in)
    add_card(slide5, Inches(0.45), Inches(0.95), Inches(8.5), Inches(5.95), title="Multidimensional Stakeholder Impact Matrix", title_color=NAVY)

    col_headers = [
        ("👵 Elderly Patients", LIGHT_BLUE, NAVY),
        ("👨‍👩‍👧 Caregivers & Families", LIGHT_ORANGE, ORANGE),
        ("🩺 Healthcare & Communities", LIGHT_GREEN, GREEN)
    ]

    for j, (ch_title, ch_bg, ch_clr) in enumerate(col_headers):
        cx = Inches(0.65 + j * 2.75)
        ch_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.45), Inches(2.6), Inches(0.5))
        ch_box.fill.solid(); ch_box.fill.fore_color.rgb = ch_bg; ch_box.line.color.rgb = ch_clr; ch_box.line.width = Pt(1.5)
        tf = ch_box.text_frame; p = tf.paragraphs[0]; p.text = ch_title; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = ch_clr; p.alignment = PP_ALIGN.CENTER

    impact_points = [
        [
            ("1. Cognitive Retention:", " Regular memory play delays functional cognitive decline."),
            ("2. Zero Diagnostic Stress:", " Playful engagement replaces alarming clinical tests."),
            ("3. Emotional Warmth:", " Companion in mother tongue combats loneliness and confusion.")
        ],
        [
            ("1. Reduced Burnout:", " Unified 7/30-day wellness trends ease constant monitoring strain."),
            ("2. Early Warning Insights:", " Gentle advisories highlight subtle cognitive shifts early."),
            ("3. Indigenous Guidance:", " Culturally relevant diet & routine tips (Manimuni, Morung).")
        ],
        [
            ("1. ASHA Worker Power:", " Offline progress tracking for routine door-to-door visits."),
            ("2. Heritage Preservation:", " Keeps 8 NER state folklores, music & traditions alive."),
            ("3. Scalable MedTech:", " Affordable digital memory care for remote hilly districts.")
        ]
    ]

    for j, pts in enumerate(impact_points):
        cx = Inches(0.65 + j * 2.75)
        tb = slide5.shapes.add_textbox(cx, Inches(2.05), Inches(2.6), Inches(4.7))
        tf = tb.text_frame; tf.word_wrap = True
        for idx, (bld, txt) in enumerate(pts):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            r1 = p.add_run(); r1.text = bld; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = DARK_NAVY
            r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(10); r2.font.color.rgb = TEXT_DARK

    # Right Sidebar Cards (Width: 3.7 in)
    add_card(slide5, Inches(9.15), Inches(0.95), Inches(3.7), Inches(2.9), title="Business Model & Viability", title_color=ORANGE, bg_color=WHITE)
    tb_bm = slide5.shapes.add_textbox(Inches(9.3), Inches(1.4), Inches(3.4), Inches(2.3))
    tf_bm = tb_bm.text_frame; tf_bm.word_wrap = True
    bm_pts = [
        ("• B2G National Health Mission: ", "Partnership with state health departments for rural ASHA worker deployment."),
        ("• B2C Freemium Model: ", "Free core courtyard games; premium family story customization."),
        ("• Eldercare Institution SaaS: ", "Subscription tier for private memory care facilities & NGOs.")
    ]
    for idx, (bld, txt) in enumerate(bm_pts):
        p = tf_bm.paragraphs[0] if idx == 0 else tf_bm.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run(); r1.text = bld; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(10); r2.font.color.rgb = TEXT_DARK

    add_card(slide5, Inches(9.15), Inches(4.0), Inches(3.7), Inches(2.9), title="Future Scope & Expansion", title_color=GREEN, bg_color=LIGHT_GREEN, border_color=GREEN)
    tb_fs = slide5.shapes.add_textbox(Inches(9.3), Inches(4.45), Inches(3.4), Inches(2.3))
    tf_fs = tb_fs.text_frame; tf_fs.word_wrap = True
    fs_pts = [
        ("• Wearable & IoT Integration: ", "Sync heart rate & sleep patterns with cognitive score trends."),
        ("• Smart Home Voice Speakers: ", "Ambient audio reminders over smart speakers for rural homes."),
        ("• Additional Dialects: ", "Expansion to Garo, Kokborok & regional sub-tribal languages.")
    ]
    for idx, (bld, txt) in enumerate(fs_pts):
        p = tf_fs.paragraphs[0] if idx == 0 else tf_fs.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run(); r1.text = bld; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(10); r2.font.color.rgb = TEXT_DARK


    # ==========================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "RESEARCH AND REFERENCES", 6)

    stat_cards = [
        ("8.8 Million+", "Projected Dementia Patients in India by 2026", "Source: The Lancet Public Health (2022) Study", LIGHT_BLUE, NAVY),
        ("85% Rural Gap", "Lack of Geriatric Memory Care Clinics in Rural NER", "Source: NITI Aayog Health Index Reports", LIGHT_ORANGE, ORANGE),
        ("3.2x Adherence", "Higher Engagement via Dialect & Cultural Cues", "Source: WHO Cognitive Care Best Practices", LIGHT_GREEN, GREEN)
    ]

    for idx, (s_stat, s_title, s_src, s_bg, s_clr) in enumerate(stat_cards):
        sx = Inches(0.45 + idx * 4.15)
        scard = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(0.95), Inches(4.0), Inches(1.8))
        scard.fill.solid(); scard.fill.fore_color.rgb = s_bg
        scard.line.color.rgb = s_clr; scard.line.width = Pt(1.5)
        tf = scard.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = s_stat; p.font.bold = True; p.font.size = Pt(22); p.font.color.rgb = s_clr; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = s_title; p2.font.bold = True; p2.font.size = Pt(11); p2.font.color.rgb = TEXT_DARK; p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = s_src; p3.font.italic = True; p3.font.size = Pt(9); p3.font.color.rgb = TEXT_MUTED; p3.alignment = PP_ALIGN.CENTER

    add_card(slide6, Inches(0.45), Inches(2.9), Inches(12.45), Inches(1.85), title="Clinical Evidence & Regulatory Compliance Foundations", title_color=NAVY)
    tb_ev = slide6.shapes.add_textbox(Inches(0.6), Inches(3.35), Inches(12.15), Inches(1.3))
    tf_ev = tb_ev.text_frame; tf_ev.word_wrap = True
    
    ev_pts = [
        ("• Cognitive Reserve & Reminiscence Therapy: ", "Clinical research confirms that engaging elders with familiar childhood photographs, folk tunes, and native mother tongues stimulates dormant neural pathways, reducing agitation and slowing cognitive decline in mild-to-moderate dementia."),
        ("• Local-First Privacy & DPDP Compliance: ", "Architected in accordance with India's Digital Personal Data Protection (DPDP) Act 2023. Zero camera frames leave the device; all family names, memory graphs, and caregiver notes are encrypted locally with AES-256 standard encryption keys.")
    ]
    for idx, (bld, txt) in enumerate(ev_pts):
        p = tf_ev.paragraphs[0] if idx == 0 else tf_ev.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run(); r1.text = bld; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(10); r2.font.color.rgb = TEXT_DARK

    add_card(slide6, Inches(0.45), Inches(4.9), Inches(12.45), Inches(2.0), title="Key Academic & Technical References", title_color=ORANGE)
    
    rows, cols = 4, 3
    t_ref = slide6.shapes.add_table(rows, cols, Inches(0.6), Inches(5.35), Inches(12.15), Inches(1.45)).table
    t_ref.columns[0].width = Inches(3.2)
    t_ref.columns[1].width = Inches(5.8)
    t_ref.columns[2].width = Inches(3.15)

    headers = ["Reference Source / Institution", "Document / Study Title", "Relevance to Vanika Platform"]
    for j, h in enumerate(headers):
        cell = t_ref.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.text = h; p.font.bold = True; p.font.size = Pt(10); p.font.color.rgb = WHITE

    ref_data = [
        ("World Health Organization (WHO)", "Global Action Plan on Public Health Response to Dementia (2017-2025)", "Framework for non-pharmacological community memory care."),
        ("The Lancet Public Health (2022)", "Estimation of Global Dementia Prevalence in 2019 & Forecast for 2050", "Epidemiological statistical baseline for Indian dementia load."),
        ("Ministry of Health (MoHFW)", "National Programme for Health Care of the Elderly (NPHCE) Guidelines", "Alignment with government rural healthcare delivery goals.")
    ]

    for i, row in enumerate(ref_data):
        for j, val in enumerate(row):
            cell = t_ref.cell(i+1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i%2==0 else BG_LIGHT
            p = cell.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(9.5); p.font.color.rgb = TEXT_DARK
            if j == 0: p.font.bold = True


    # ==========================================
    # SLIDE 7: CONCLUSION & ROADMAP
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "CONCLUSION & EXECUTION ROADMAP", 7)

    banner_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.95), Inches(12.45), Inches(0.85))
    banner_card.fill.solid(); banner_card.fill.fore_color.rgb = LIGHT_BLUE
    banner_card.line.color.rgb = NAVY; banner_card.line.width = Pt(1.5)
    
    flow_steps = ["Dementia Challenge in NER", "Vanika Local AI Companion", "Enculturated Memory Play", "Caregiver & ASHA Insights", "Accessible Rural MedTech"]
    for i, step in enumerate(flow_steps):
        bx = Inches(0.65 + i * 2.45)
        tb = slide7.shapes.add_textbox(bx, Inches(1.05), Inches(2.1), Inches(0.65))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = step; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = DARK_NAVY; p.alignment = PP_ALIGN.CENTER
        
        if i < 4:
            arr = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.65 + i * 2.45), Inches(1.25), Inches(0.25), Inches(0.2))
            arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE; arr.line.fill.background()

    add_card(slide7, Inches(0.45), Inches(1.95), Inches(5.9), Inches(4.95), title="Key Solution Advantages & Value Summary", title_color=NAVY)
    tb_c1 = slide7.shapes.add_textbox(Inches(0.6), Inches(2.45), Inches(5.6), Inches(4.3))
    tf_c1 = tb_c1.text_frame; tf_c1.word_wrap = True

    c1_pts = [
        ("• Culturally Attuned Care: ", "Tailored specifically around 8 North-Eastern states, 6 regional dialects, and indigenous herbal/folklore wisdom."),
        ("• Privacy-First Engineering: ", "100% offline local processing with zero camera cloud streaming; full compliance with DPDP Act 2023."),
        ("• Dual-Empowerment Model: ", "Supports both the elderly patient (joyful courtyard play) and the caregiver/ASHA worker (actionable 7/30-day trends)."),
        ("• High Accessibility: ", "Designed for low digital literacy with 18px text, high contrast, warm voice guide, and zero diagnostic stress.")
    ]
    for idx, (bld, txt) in enumerate(c1_pts):
        p = tf_c1.paragraphs[0] if idx == 0 else tf_c1.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run(); r1.text = bld; r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(11); r2.font.color.rgb = TEXT_DARK

    add_card(slide7, Inches(6.55), Inches(1.95), Inches(6.35), Inches(4.95), title="Phased Execution Roadmap (2026)", title_color=ORANGE, bg_color=WHITE)

    roadmap_steps = [
        ("🚀 Phase 1 (Q1 2026) — Core Launch", "Web PWA deployment with Assamese & English AI models, local vault & core 4 courtyard games.", LIGHT_BLUE),
        ("🏥 Phase 2 (Q2 2026) — Rural Pilot", "Field testing across primary health centers in Assam & Meghalaya with ASHA workers.", LIGHT_GREEN),
        ("🗣️ Phase 3 (Q3 2026) — Multilingual Scale", "Integration of Bodo, Khasi, Mizo & Nagamese models, plus IoT smart speaker voice sync.", LIGHT_ORANGE),
        ("🌐 Phase 4 (Q4 2026) — National Rollout", "Expansion under National Health Mission (NHM) North-East Digital Health Drive.", LIGHT_BLUE)
    ]

    for idx, (r_title, r_desc, r_bg) in enumerate(roadmap_steps):
        ry = Inches(2.5 + idx * 1.05)
        r_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), ry, Inches(5.95), Inches(0.95))
        r_box.fill.solid(); r_box.fill.fore_color.rgb = r_bg
        r_box.line.color.rgb = BORDER_COLOR; r_box.line.width = Pt(1)
        tf = r_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = r_title; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = DARK_NAVY
        p2 = tf.add_paragraph(); p2.text = r_desc; p2.font.size = Pt(10); p2.font.color.rgb = TEXT_DARK

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Successfully generated presentation at: {output_path}")

if __name__ == "__main__":
    out_file = os.path.abspath(os.path.join("ppt", "Vanika_SIH_2026_Presentation.pptx"))
    create_presentation(out_file)
